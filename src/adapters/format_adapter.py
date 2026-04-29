from abc import ABC, abstractmethod
import re
import os

# 确保输出目录存在
os.makedirs("output", exist_ok=True)


class Exporter(ABC):
    """导出器抽象基类"""
    
    @abstractmethod
    def export(self, markdown_content: str, filename: str) -> str:
        """
        将Markdown内容导出为指定格式
        :param markdown_content: Markdown格式内容
        :param filename: 输出文件名（不含扩展名）
        :return: 生成的文件路径
        """
        pass
    
    @abstractmethod
    def get_extension(self) -> str:
        """返回导出格式的扩展名"""
        pass


class MarkdownExporter(Exporter):
    """Markdown格式导出器"""
    
    def export(self, markdown_content: str, filename: str) -> str:
        filepath = f"output/{filename}.md"
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        return filepath
    
    def get_extension(self) -> str:
        return '.md'


class TxtExporter(Exporter):
    """纯文本格式导出器"""
    
    def export(self, markdown_content: str, filename: str) -> str:
        text = re.sub(r'[#*`>\-\[\]]', '', markdown_content)
        text = re.sub(r'\n{2,}', '\n\n', text).strip()
        
        filepath = f"output/{filename}.txt"
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(text)
        return filepath
    
    def get_extension(self) -> str:
        return '.txt'


class PptExporter(Exporter):
    """PPT格式导出器 - 增强版"""

    def export(self, markdown_content: str, filename: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("请先安装 python-pptx: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        lines = markdown_content.split('\n')
        current_slide = None
        current_content = []
        in_code_block = False
        code_block_content = []

        def add_content_slide(prs, title_text, content_lines):
            if not content_lines:
                return
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if title_text:
                title = slide.shapes.title
                title.text = title_text
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    if shape != slide.shapes.title:
                        tf = shape.text_frame
                        tf.clear()
                        for i, line in enumerate(content_lines):
                            if i == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = line
                            p.level = 0
                        break
            return slide

        for line in lines:
            if line.strip().startswith('```'):
                if not in_code_block:
                    in_code_block = True
                    code_block_content = []
                else:
                    in_code_block = False
                    if current_slide is None:
                        current_slide = add_content_slide(prs, "代码", code_block_content)
                    else:
                        current_content.extend(['```'] + code_block_content + ['```'])
                    code_block_content = []
                continue

            if in_code_block:
                code_block_content.append(line)
                continue

            if line.startswith('# '):
                if current_content and current_slide:
                    add_content_slide(prs, "", current_content)
                    current_content = []
                title = line[2:].strip()
                slide_layout = prs.slide_layouts[0]
                current_slide = prs.slides.add_slide(slide_layout)
                current_slide.shapes.title.text = title
                current_content = []

            elif line.startswith('## '):
                if current_content:
                    add_content_slide(prs, "", current_content)
                    current_content = []
                title = line[3:].strip()
                slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(slide_layout)
                current_slide.shapes.title.text = title

            elif line.startswith('### '):
                if current_content:
                    add_content_slide(prs, "", current_content)
                    current_content = []
                title = line[4:].strip()
                slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(slide_layout)
                current_slide.shapes.title.text = title

            elif line.startswith('- ') or line.startswith('* '):
                current_content.append('• ' + line[2:].strip())

            elif line.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
                match = re.match(r'^(\d+)\. (.+)$', line)
                if match:
                    current_content.append(f"{match.group(1)}. {match.group(2)}")

            elif line.strip().startswith('|'):
                continue

            elif line.strip() and not line.startswith('#'):
                stripped = line.strip()
                if stripped and stripped not in ['```', '---', '***']:
                    current_content.append(stripped)

        if current_content:
            add_content_slide(prs, "", current_content)

        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "文档"

        filepath = f"output/{filename}.pptx"
        prs.save(filepath)
        return filepath

    def get_extension(self) -> str:
        return '.pptx'


class FormatAdapter:
    """格式转换适配器"""
    
    def __init__(self):
        self._exporters = {
            'md': MarkdownExporter(),
            'txt': TxtExporter(),
            'ppt': PptExporter()
        }
    
    def export(self, markdown_content: str, filename: str, format_type: str) -> str:
        """
        统一导出接口
        :param markdown_content: Markdown内容
        :param filename: 文件名
        :param format_type: 输出格式（md/txt/ppt）
        :return: 文件路径
        """
        if format_type not in self._exporters:
            raise ValueError(f"不支持的格式: {format_type}")
        
        exporter = self._exporters[format_type]
        return exporter.export(markdown_content, filename)
    
    def get_supported_formats(self) -> list:
        """返回支持的格式列表"""
        return list(self._exporters.keys())